# -*- coding: utf-8 -*-
"""Презентация LUIS+ · SEO & GEO · Июнь 2026 — в фирменном стиле LUIS+
Дизайн взят из шаблона LUIS+ 2026_расширенная.pptx:
  красный #BE1730, сталь #7D98AB, тёмный #1A1A1A, панель #ECF0F8, жёлтый #FFC72C,
  шрифты Exo 2 SemiBold / Exo 2, белые слайды, логотип слева сверху."""
import colorsys, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# ---------- ДАННЫЕ (из luis-dashboard.html) ----------
CATS = ["Комплексные","Видеонаблюдение","ОПС","Пожаротушение","СКУД","Оповещение","Электропитание","Кабель/электрика"]
CAT_N = {"Комплексные":18,"Видеонаблюдение":17,"ОПС":50,"Пожаротушение":34,"СКУД":28,"Оповещение":13,"Электропитание":12,"Кабель/электрика":91}
BRANDS = ["LUIS+","etm.ru","rsvet.ru","rubezh.ru","dssl.ru","bolid.ru","tinko.ru","layta.ru","satro-paladin.com","rutektd.ru","iss.ru","dean.ru","mossb.ru","garantgroup.com","argus-spectr.ru"]

SEO_M = {"Комплексные":{"LUIS+":6,"etm.ru":0,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":6,"bolid.ru":0,"tinko.ru":0,"layta.ru":0,"satro-paladin.com":0,"rutektd.ru":0,"iss.ru":11,"dean.ru":0,"mossb.ru":6,"garantgroup.com":0,"argus-spectr.ru":0},
"Видеонаблюдение":{"LUIS+":0,"etm.ru":0,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":24,"bolid.ru":0,"tinko.ru":29,"layta.ru":0,"satro-paladin.com":6,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":0},
"ОПС":{"LUIS+":18,"etm.ru":36,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":0,"bolid.ru":12,"tinko.ru":90,"layta.ru":8,"satro-paladin.com":72,"rutektd.ru":0,"iss.ru":0,"dean.ru":8,"mossb.ru":0,"garantgroup.com":2,"argus-spectr.ru":4},
"Пожаротушение":{"LUIS+":26,"etm.ru":24,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":0,"bolid.ru":0,"tinko.ru":62,"layta.ru":0,"satro-paladin.com":56,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":0},
"СКУД":{"LUIS+":0,"etm.ru":57,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":36,"bolid.ru":4,"tinko.ru":75,"layta.ru":0,"satro-paladin.com":43,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":0},
"Оповещение":{"LUIS+":0,"etm.ru":8,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":0,"bolid.ru":15,"tinko.ru":54,"layta.ru":0,"satro-paladin.com":0,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":0},
"Электропитание":{"LUIS+":8,"etm.ru":33,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":0,"bolid.ru":0,"tinko.ru":100,"layta.ru":8,"satro-paladin.com":75,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":0},
"Кабель/электрика":{"LUIS+":0,"etm.ru":90,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":0,"bolid.ru":0,"tinko.ru":54,"layta.ru":9,"satro-paladin.com":27,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":0}}

GEO_M = {"Комплексные":{"LUIS+":72,"etm.ru":0,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":0,"bolid.ru":0,"tinko.ru":0,"layta.ru":0,"satro-paladin.com":0,"rutektd.ru":0,"iss.ru":0,"dean.ru":6,"mossb.ru":0,"garantgroup.com":6,"argus-spectr.ru":0},
"Видеонаблюдение":{"LUIS+":12,"etm.ru":0,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":18,"bolid.ru":0,"tinko.ru":6,"layta.ru":0,"satro-paladin.com":0,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":0},
"ОПС":{"LUIS+":6,"etm.ru":4,"rsvet.ru":2,"rubezh.ru":22,"dssl.ru":0,"bolid.ru":48,"tinko.ru":56,"layta.ru":18,"satro-paladin.com":42,"rutektd.ru":0,"iss.ru":0,"dean.ru":16,"mossb.ru":2,"garantgroup.com":22,"argus-spectr.ru":26},
"Пожаротушение":{"LUIS+":18,"etm.ru":3,"rsvet.ru":0,"rubezh.ru":3,"dssl.ru":0,"bolid.ru":15,"tinko.ru":9,"layta.ru":6,"satro-paladin.com":21,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":9,"argus-spectr.ru":0},
"СКУД":{"LUIS+":100,"etm.ru":4,"rsvet.ru":0,"rubezh.ru":11,"dssl.ru":50,"bolid.ru":4,"tinko.ru":21,"layta.ru":7,"satro-paladin.com":21,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":0},
"Оповещение":{"LUIS+":8,"etm.ru":15,"rsvet.ru":0,"rubezh.ru":54,"dssl.ru":0,"bolid.ru":0,"tinko.ru":54,"layta.ru":38,"satro-paladin.com":8,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":0,"garantgroup.com":8,"argus-spectr.ru":0},
"Электропитание":{"LUIS+":92,"etm.ru":25,"rsvet.ru":0,"rubezh.ru":8,"dssl.ru":17,"bolid.ru":17,"tinko.ru":67,"layta.ru":0,"satro-paladin.com":42,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":8,"garantgroup.com":0,"argus-spectr.ru":0},
"Кабель/электрика":{"LUIS+":0,"etm.ru":12,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":0,"bolid.ru":0,"tinko.ru":9,"layta.ru":7,"satro-paladin.com":16,"rutektd.ru":0,"iss.ru":0,"dean.ru":0,"mossb.ru":1,"garantgroup.com":0,"argus-spectr.ru":0}}

SEO_O = {"LUIS+":8,"etm.ru":49,"rsvet.ru":0,"rubezh.ru":0,"dssl.ru":6,"bolid.ru":3,"tinko.ru":61,"layta.ru":5,"satro-paladin.com":39,"rutektd.ru":0,"iss.ru":1,"dean.ru":2,"mossb.ru":0,"garantgroup.com":0,"argus-spectr.ru":1}
GEO_O = {"LUIS+":24,"etm.ru":8,"rsvet.ru":0,"rubezh.ru":9,"dssl.ru":7,"bolid.ru":12,"tinko.ru":23,"layta.ru":9,"satro-paladin.com":21,"rutektd.ru":0,"iss.ru":0,"dean.ru":3,"mossb.ru":1,"garantgroup.com":6,"argus-spectr.ru":5}

def short(b): return b.replace('.ru','').replace('.com','')

# ---------- ФИРМЕННЫЕ ЦВЕТА LUIS+ ----------
CR    = RGBColor(0xBE,0x17,0x30)   # красный LUIS+
CRD   = RGBColor(0x8E,0x10,0x24)   # тёмно-красный
STEEL = RGBColor(0x7D,0x98,0xAB)   # сталь
STEELD= RGBColor(0x0F,0x4C,0x81)   # тёмно-синий
INK   = RGBColor(0x1A,0x1A,0x1A)   # тёмный текст
INK2  = RGBColor(0x44,0x4A,0x52)
GRAY  = RGBColor(0x8A,0x93,0x9E)
PANEL = RGBColor(0xEC,0xF0,0xF8)   # светлая панель
PANEL2= RGBColor(0xF5,0xF7,0xFB)
LINE  = RGBColor(0xD9,0xE0,0xEA)
WHITE = RGBColor(0xFF,0xFF,0xFF)
GOOD  = RGBColor(0x1f,0x8a,0x52)
YELLOW= RGBColor(0xFF,0xC7,0x2C)
BADBG = RGBColor(0xFB,0xE7,0xEA)

HEAD = "Exo 2 SemiBold"
BODY = "Exo 2"
LOGO = "/Users/daria/Desktop/Сlaude Code/LUIS+/luis_logo.png"
TEMPLATE = "/Users/daria/Desktop/Сlaude Code/LUIS+/Презентации/LUIS+ 2026_расширенная.pptx"

def heat(p):
    """Тепловая шкала в фирменной сине-стальной гамме (светлая → тёмно-синяя)."""
    if p == 0:
        return RGBColor(0xF3,0xF5,0xF9), RGBColor(0xBD,0xC5,0xD0)
    t = (p/100.0)**0.82
    lo = (0xDD,0xE7,0xF2); hi = (0x0F,0x4C,0x81)
    r = int(lo[0]+(hi[0]-lo[0])*t); g = int(lo[1]+(hi[1]-lo[1])*t); b = int(lo[2]+(hi[2]-lo[2])*t)
    txt = WHITE if t > 0.45 else STEELD
    return RGBColor(r,g,b), txt

# ---------- КАРКАС: основа — фирменный шаблон LUIS+ ----------
prs = Presentation(TEMPLATE)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_masters[0].slide_layouts[1]  # «Пустой без логотипа»

def bgfill(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color

# оставляем только родной титульный слайд, остальные удаляем
def trim_to_title():
    lst = prs.slides._sldIdLst
    for sldId in list(lst)[1:]:
        rId = sldId.get(qn('r:id'))
        try: prs.part.drop_rel(rId)
        except Exception: pass
        lst.remove(sldId)

# заменить текст абзацев в текстовом кадре, сохранив форматирование
def set_para_texts(tf, texts):
    for i, para in enumerate(tf.paragraphs):
        if i < len(texts):
            if para.runs:
                para.runs[0].text = texts[i]
                for r in para.runs[1:]: r.text = ""
            else:
                r = para.add_run(); r.text = texts[i]
        else:
            for r in para.runs: r.text = ""

# адаптировать заголовок родного титульника под тему отчёта
def retitle():
    s = prs.slides[0]
    boxes = {}
    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == 6: walk(sh.shapes)
            elif sh.has_text_frame: boxes[sh.name] = sh.text_frame
    walk(s.shapes)
    if "TextBox 4" in boxes:
        set_para_texts(boxes["TextBox 4"], ["ВИДИМОСТЬ LUIS+", "В ПОИСКЕ", "И НЕЙРОСЕТЯХ"])
    if "TextBox 5" in boxes:
        set_para_texts(boxes["TextBox 5"], ["SEO & GEO · аналитика видимости · июнь 2026", "263 запроса · LUIS+ и 14 конкурентов"])

trim_to_title()
retitle()

# создать чистый слайд на пустом макете
def blank_slide():
    s = prs.slides.add_slide(BLANK)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    bgfill(s, WHITE)
    return s

def rect(s, x, y, w, h, fill=None, line=None, line_w=None, radius=False, shadow=False, shape=None):
    shp = s.shapes.add_shape(shape or (MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE), x, y, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {}); el.append(ef)
        sh = ef.makeelement(qn('a:outerShdw'), {'blurRad':'90000','dist':'38000','dir':'5400000','rotWithShape':'0'}); ef.append(sh)
        clr = sh.makeelement(qn('a:srgbClr'), {'val':'1A2A3A'}); sh.append(clr)
        clr.append(clr.makeelement(qn('a:alpha'), {'val':'11000'}))
    return shp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=2, line_sp=1.0, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0); p.line_spacing = line_sp
        for (t,sz,b,c,*rest) in para:
            fam = rest[1] if len(rest)>1 else (HEAD if b else BODY)
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = False; r.font.name = fam
            r.font.color.rgb = c; r.font.italic = it
    return tb

def logo(s, x=Inches(0.62), y=Inches(0.42), w=Inches(1.55)):
    s.shapes.add_picture(LOGO, x, y, width=w)  # ratio 598:203

def head(s, kicker, title, sub=None):
    """Шапка слайда в стиле LUIS+: логотип · красный кикер · тёмный заголовок · красная черта."""
    logo(s)
    txt(s, Inches(0.64), Inches(1.16), Inches(11.5), Inches(0.3),
        [[(kicker, 12.5, True, CR, False, HEAD)]])
    txt(s, Inches(0.62), Inches(1.46), Inches(12.1), Inches(0.6),
        [[(title, 25, True, INK, False, HEAD)]])
    rect(s, Inches(0.64), Inches(2.06), Inches(0.62), Pt(3.4), fill=CR)
    if sub:
        txt(s, Inches(0.64), Inches(2.18), Inches(12), Inches(0.35), [[(sub, 12, False, STEEL, False, BODY)]])

def footer(s):
    txt(s, Inches(0.62), Inches(7.08), Inches(9), Inches(0.3),
        [[("Источник: вкладка «Июнь 2026» · 263 запроса · LUIS+ и 14 конкурентов", 8.5, False, GRAY, False, BODY)]])
    txt(s, Inches(9.7), Inches(7.08), Inches(3.03), Inches(0.3),
        [[("LUIS+ · Отдел маркетинга · Июнь 2026", 8.5, False, GRAY, False, BODY)]], align=PP_ALIGN.RIGHT)

# ================= СЛАЙД 1 — родной титульник шаблона (уже в prs) =================

# ================= СЛАЙД 2 — KPI + ВЫВОД =================
s = blank_slide()
head(s, "КЛЮЧЕВЫЕ ПОКАЗАТЕЛИ", "Итоги видимости за июнь 2026")
KPI = [
    ("СЕМАНТИКА", "263", STEEL, "987 567 показов/мес · 8 групп", None),
    ("GEO · ТОП-10 В ИИ", "24%", CR, "лучше конкурентов, но уровень низкий · 64 запроса", ("наращивать", CR, BADBG)),
    ("SEO · ТОП-10 ЯНДЕКСА", "8%", CR, "низкая база · медиана позиции 28 · охват топ-100 44%", ("слабо", CR, BADBG)),
]
cx, cy, cw, ch, gap = Inches(0.62), Inches(2.55), Inches(3.87), Inches(1.95), Inches(0.235)
for i,(lab,num,ncol,ctx,tag) in enumerate(KPI):
    x = cx + i*(cw+gap)
    rect(s, x, cy, cw, ch, fill=WHITE, line=LINE, line_w=Pt(1), radius=True, shadow=True)
    rect(s, x, cy, Pt(4), ch, fill=CR)
    txt(s, x+Inches(0.28), cy+Inches(0.2), cw-Inches(0.45), Inches(0.35), [[(lab, 10.5, True, GRAY, False, HEAD)]])
    txt(s, x+Inches(0.26), cy+Inches(0.46), cw-Inches(0.45), Inches(0.7), [[(num, 42, True, ncol, False, HEAD)]])
    if tag:
        rect(s, x+Inches(0.27), cy+Inches(1.22), Inches(1.2), Inches(0.3), fill=tag[2], radius=True)
        txt(s, x+Inches(0.27), cy+Inches(1.25), Inches(1.2), Inches(0.26), [[(tag[0], 9.5, True, tag[1], False, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.28), cy+Inches(1.58), cw-Inches(0.5), Inches(0.4), [[(ctx, 10.5, False, INK2, False, BODY)]], line_sp=1.06)
# врезка-вывод
by = Inches(4.85)
rect(s, Inches(0.62), by, Inches(12.1), Inches(1.95), fill=PANEL, line=None, radius=True)
rect(s, Inches(0.62), by, Pt(5), Inches(1.95), fill=CR)
txt(s, Inches(0.95), by+Inches(0.24), Inches(11.5), Inches(0.5),
    [[("Главный вывод: ", 16, True, CR, False, HEAD),("в нейросетях (GEO) мы впереди конкурентов, но охват всё равно низкий — 24%. В классическом SEO — совсем низко, 8%. Обе зоны требуют роста.", 16, False, INK, False, BODY)]], line_sp=1.08)
bullets = [
    "GEO (нейросети): опережаем всех 14 конкурентов — но 24% в топ-10 это стартовая база, а не победа.",
    "SEO (Яндекс): отстаём от лидеров — 8% топ-10 против 61% у tinko.ru и 49% у etm.ru, медиана позиции 28.",
    "Задача — растить обе зоны: точки силы в ИИ (СКУД, Электропитание, Комплексные) переносить в органику Яндекса.",
]
for i,b in enumerate(bullets):
    yb = by+Inches(0.82)+i*Inches(0.35)
    rect(s, Inches(0.98), yb+Inches(0.05), Inches(0.1), Inches(0.1), fill=CR, shape=MSO_SHAPE.DIAMOND)
    txt(s, Inches(1.25), yb, Inches(11.2), Inches(0.35), [[(b, 12, False, INK2, False, BODY)]])
footer(s)

# ================= СЛАЙД 3 — ГЕРОЙ GEO =================
s = blank_slide()
head(s, "GEO · ТОП-10 ОТВЕТОВ CHATGPT · GEMINI · PERPLEXITY", "В нейросетях опережаем конкурентов")
txt(s, Inches(0.64), Inches(2.24), Inches(12), Inches(0.4),
    [[("Опережаем всех 14 конкурентов по видимости в ИИ, но 24% в топ-10 — стартовая база: приоритет — наращивать присутствие.", 12.5, False, INK2, False, BODY)]])
top = sorted(GEO_O.items(), key=lambda kv:-kv[1])[:5]
mx = max(v for _,v in top) or 1
lx, ly, lw = Inches(0.62), Inches(2.85), Inches(7.6)
rect(s, lx, ly, lw, Inches(3.75), fill=PANEL, radius=True)
txt(s, lx+Inches(0.35), ly+Inches(0.24), lw-Inches(0.6), Inches(0.3),
    [[("ДОЛЯ ЗАПРОСОВ В ТОП-10 ИИ-ОТВЕТОВ", 11, True, STEELD, False, HEAD)]])
rowY = ly+Inches(0.7); rowH = Inches(0.4); rowGap = Inches(0.17)
barX = lx+Inches(2.3); barMaxW = Inches(4.0)
for i,(b,v) in enumerate(top):
    y = rowY + i*(rowH+rowGap); me = (b=="LUIS+")
    txt(s, lx+Inches(0.35), y+Inches(0.02), Inches(0.3), rowH, [[(str(i+1), 14, True, STEEL, False, HEAD)]])
    txt(s, lx+Inches(0.66), y+Inches(0.02), Inches(1.7), rowH, [[(short(b), 13, True, CR if me else INK, False, HEAD)]])
    rect(s, barX, y+Inches(0.08), barMaxW, Inches(0.24), fill=WHITE, radius=True)
    fw = int(barMaxW*(v/mx))
    if fw>0: rect(s, barX, y+Inches(0.08), Emu(fw), Inches(0.24), fill=CR if me else STEEL, radius=True)
    txt(s, barX+barMaxW+Inches(0.12), y, Inches(0.85), rowH, [[(str(v)+"%", 15, True, CR if me else INK, False, HEAD)]], align=PP_ALIGN.RIGHT)
# чипы справа
chips = [("64","запроса с LUIS+ в ИИ"),("55","с разбросом по движкам"),("3","движка: ChatGPT · Gemini · Perplexity")]
chx, chy, chw, chh = Inches(8.42), Inches(2.85), Inches(4.3), Inches(1.13)
for i,(cv,cl) in enumerate(chips):
    y = chy + i*(chh+Inches(0.18))
    rect(s, chx, y, chw, chh, fill=WHITE, line=LINE, line_w=Pt(1), radius=True, shadow=True)
    rect(s, chx, y, Pt(4), chh, fill=STEEL)
    txt(s, chx+Inches(0.3), y+Inches(0.16), chw-Inches(0.5), Inches(0.6), [[(cv, 30, True, CR, False, HEAD)]])
    txt(s, chx+Inches(0.3), y+Inches(0.74), chw-Inches(0.5), Inches(0.35), [[(cl, 12, False, INK2, False, BODY)]])
footer(s)

# ================= СЛАЙДЫ 4–5 — ЛИДЕРЫ (бары) =================
def leader_bar_slide(kicker, title, sub, data, note):
    s = blank_slide()
    head(s, kicker, title, sub)
    items = sorted([(k,v) for k,v in data.items() if v>0], key=lambda kv:-kv[1])
    cats = [short(k) for k,_ in items]; vals = [v for _,v in items]
    cd = CategoryChartData(); cd.categories = cats; cd.add_series("% в топ-10", vals)
    gx, gy, gw, gh = Inches(0.62), Inches(2.55), Inches(12.1), Inches(4.05)
    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, gx, gy, gw, gh, cd)
    ch = gf.chart; ch.has_legend = False; ch.has_title = False
    plot = ch.plots[0]; plot.gap_width = 55; ser = plot.series[0]
    for idx, cat in enumerate(cats):
        pt = ser.points[idx]; pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = CR if cat=="LUIS+" else STEEL
    ser.data_labels.show_value = True
    ser.data_labels.number_format = '0"%"'; ser.data_labels.number_format_is_linked = False
    ser.data_labels.font.size = Pt(11); ser.data_labels.font.bold = True; ser.data_labels.font.name = HEAD
    ser.data_labels.font.color.rgb = INK2
    ca = ch.category_axis; ca.tick_labels.font.size = Pt(12); ca.tick_labels.font.name = HEAD; ca.tick_labels.font.bold = False
    va = ch.value_axis; va.has_major_gridlines = True
    va.tick_labels.font.size = Pt(10); va.tick_labels.number_format = '0"%"'; va.tick_labels.number_format_is_linked = False
    va.maximum_scale = max(vals)+8
    txt(s, Inches(0.62), Inches(6.62), Inches(12.1), Inches(0.4), [[(note, 11.5, False, INK2, True, BODY)]])
    footer(s)

leader_bar_slide("ИТОГОВЫЕ ЛИДЕРЫ · GEO", "GEO · Нейросети: доля запросов в топ-10 ИИ-ответов",
    "По всем 263 запросам · ChatGPT · Gemini · Perplexity", GEO_O,
    "LUIS+ (24%) впереди конкурентов, но отрыв от tinko.ru (23%) и satro-paladin (21%) минимальный, а абсолютный уровень низкий. ИИ-видимость нужно наращивать.")

leader_bar_slide("ИТОГОВЫЕ ЛИДЕРЫ · SEO", "SEO · Яндекс: доля запросов в топ-10 органики",
    "По всем 263 запросам · органическая выдача Яндекса", SEO_O,
    "SEO — главная зона роста: 8% топ-10 против 61% у tinko.ru и 49% у etm.ru. Разрыв с лидерами кратный.")

# ================= СЛАЙДЫ 6–7 — МАТРИЦЫ =================
def matrix_slide(kicker, title, sub, M, O, note):
    s = blank_slide()
    head(s, kicker, title, sub)
    rows = len(CATS)+2; cols = len(BRANDS)+1
    tx, ty = Inches(0.34), Inches(2.55); tw, th = Inches(12.66), Inches(3.95)
    gf = s.shapes.add_table(rows, cols, tx, ty, tw, th); tbl = gf.table
    tbl.columns[0].width = Inches(1.75)
    bw = int((tw - Inches(1.75)) / len(BRANDS))
    for j in range(1, cols): tbl.columns[j].width = Emu(bw)
    tbl.first_row = False; tbl.horz_banding = False
    def setcell(r,c,text,bg,fg,size=8.5,align=PP_ALIGN.CENTER,fam=HEAD):
        cell = tbl.cell(r,c); cell.fill.solid(); cell.fill.fore_color.rgb = bg
        cell.margin_left=Emu(20000); cell.margin_right=Emu(20000); cell.margin_top=Emu(8000); cell.margin_bottom=Emu(8000)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = align
        r0 = p.add_run(); r0.text = text; r0.font.size = Pt(size); r0.font.bold = False; r0.font.name = fam; r0.font.color.rgb = fg
    setcell(0,0,"Раздел", STEELD, WHITE, 9, PP_ALIGN.LEFT)
    for j,b in enumerate(BRANDS):
        setcell(0,j+1, short(b), CR if b=="LUIS+" else STEELD, WHITE, 7.5)
    for i,cat in enumerate(CATS):
        setcell(i+1,0, cat, WHITE, INK, 8.5, PP_ALIGN.LEFT)
        for j,b in enumerate(BRANDS):
            p = M[cat][b]; bg,fg = heat(p)
            setcell(i+1,j+1, (str(p)+"%" if p else "·"), bg, fg, 8.5)
    setcell(rows-1,0, "ИТОГО (263)", PANEL, INK, 8.5, PP_ALIGN.LEFT)
    for j,b in enumerate(BRANDS):
        p = O[b]; bg,fg = heat(p); setcell(rows-1,j+1, str(p)+"%", bg, fg, 8.5)
    txt(s, Inches(0.34), Inches(6.62), Inches(12.4), Inches(0.4), [[(note, 11, False, INK2, True, BODY)]])
    footer(s)

matrix_slide("МАТРИЦА · SEO", "SEO · Яндекс: попадание в ТОП-10 по разделам",
    "% запросов раздела в топ-10 органики · темнее — выше доля · колонка LUIS+ выделена", SEO_M, SEO_O,
    "SEO: показатели низкие. Относительно сильнее в Пожаротушении (26%) и ОПС (18%); провал — Видеонаблюдение, СКУД, Оповещение, Кабель (0%). Лидеры — tinko и etm.")

matrix_slide("МАТРИЦА · GEO", "GEO · Нейросети: попадание в ТОП-10 по разделам",
    "% запросов раздела в топ-10 ИИ-ответов · темнее — выше доля · колонка LUIS+ выделена", GEO_M, GEO_O,
    "GEO: LUIS+ впереди конкурентов за счёт СКУД (100%), Электропитания (92%) и Комплексных (72%). Но общий охват ИИ — 24%: масштаб пока недостаточный.")

# ================= СЛАЙД 8 — ВЫВОДЫ =================
s = blank_slide()
head(s, "СТРАТЕГИЯ", "Выводы и приоритеты")
R = [
    ("Наращивать видимость в ИИ","В GEO мы пока лучше конкурентов (24% топ-10), но это низкая база, а не победа. ИИ — растущий канал: закреплять контентом (обзоры, сравнения, экспертные материалы) и увеличивать долю, пока не оторвались tinko.ru и satro-paladin."),
    ("SEO — главная зона роста","Топ-10 лишь 8% против 61% у tinko.ru — показатель низкий. Приоритет — поднять из «глубины» (медиана 28) ВЧ-запросы там, где в GEO уже сильны: СКУД, Электропитание, Комплексные."),
    ("Спасать Видеонаблюдение и Оповещение","0% в SEO-топ-10 и слабый GEO. Категории отданы dssl.ru и rubezh.ru. Нужны продуктовые лендинги и проработка марок."),
    ("Решить судьбу «Кабель/электрика»","91 запрос (треть семантики), но 0% топ-10 и в SEO, и в GEO — рынок держат etm.ru (90%) и DKC/IEK. Либо целевая атака, либо фокус на профильную безопасность."),
]
cw2, ch2, gapx, gapy = Inches(6.0), Inches(2.02), Inches(0.35), Inches(0.28)
x0, y0 = Inches(0.62), Inches(2.5)
for i,(h,p) in enumerate(R):
    r,c = divmod(i,2); x = x0 + c*(cw2+gapx); y = y0 + r*(ch2+gapy)
    rect(s, x, y, cw2, ch2, fill=WHITE, line=LINE, line_w=Pt(1), radius=True, shadow=True)
    rect(s, x+Inches(0.28), y+Inches(0.28), Inches(0.55), Inches(0.55), fill=CR, shape=MSO_SHAPE.HEXAGON)
    txt(s, x+Inches(0.28), y+Inches(0.35), Inches(0.55), Inches(0.45), [[(str(i+1), 18, True, WHITE, False, HEAD)]], align=PP_ALIGN.CENTER)
    txt(s, x+Inches(1.05), y+Inches(0.26), cw2-Inches(1.3), Inches(0.5), [[(h, 15.5, True, INK, False, HEAD)]], line_sp=1.0)
    txt(s, x+Inches(1.05), y+Inches(0.78), cw2-Inches(1.3), Inches(1.15), [[(p, 11.5, False, INK2, False, BODY)]], line_sp=1.1)
footer(s)

out = "/Users/daria/Desktop/Сlaude Code/LUIS+/LUIS_SEO_GEO_Июнь_2026.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("OK:", out)
print("Слайдов:", len(prs.slides._sldIdLst))
